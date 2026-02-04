#!/usr/bin/env python3
"""
MEDD SIM Pitch Deck Generator
Creates a PowerPoint presentation matching Anthony's original PDF structure
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Colors
MEDD_GREEN = RgbColor(0x0D, 0x6B, 0x56)
MEDD_LIME = RgbColor(0xC7, 0xF4, 0x64)
MEDD_DARK = RgbColor(0x3D, 0x3D, 0x3D)
WHITE = RgbColor(0xFF, 0xFF, 0xFF)
BLACK = RgbColor(0x00, 0x00, 0x00)

# Create presentation (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Get blank layout
blank_layout = prs.slide_layouts[6]

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    """Add a text box to the slide"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox

def set_background(slide, color):
    """Set slide background color"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_with_text(slide, shape_type, left, top, width, height, text, fill_color, text_color=BLACK, font_size=12, bold=False):
    """Add a shape with text"""
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT
    return shape

# Slide 1: Hero (Green background with logo)
slide1 = prs.slides.add_slide(blank_layout)
set_background(slide1, MEDD_LIME)

add_text_box(slide1, 2, 2, 9, 1.5, "medd sim", font_size=72, bold=True, color=MEDD_GREEN, align=PP_ALIGN.CENTER)
add_text_box(slide1, 1.5, 4.5, 10, 2, 
    "Giving your team the tools to practice the moments that matter,\nbefore they matter.",
    font_size=24, color=MEDD_GREEN, align=PP_ALIGN.CENTER)

# Slide 2: Value Proposition (Dark centered)
slide2 = prs.slides.add_slide(blank_layout)
set_background(slide2, MEDD_DARK)

add_text_box(slide2, 1, 2, 11, 4,
    "MEDD Sim is the simulation studio you control.\n\n"
    "Build any AI-powered coach, role-play, examiner or supporting asset in minutes—"
    "then stitch them into learning pathways that turn high-stakes conversations into "
    "rehearsed performances.",
    font_size=28, color=MEDD_LIME, align=PP_ALIGN.CENTER)

# Slide 3: Vision (Split with image placeholder)
slide3 = prs.slides.add_slide(blank_layout)
set_background(slide3, MEDD_DARK)

add_text_box(slide3, 0.5, 0.5, 2, 0.5, "Our vision", font_size=20, color=MEDD_LIME)
add_text_box(slide3, 0.5, 1.2, 5.5, 2, 
    "is to make deliberate practice a daily norm.",
    font_size=36, bold=True, color=WHITE)

# Add image placeholder
shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1), Inches(5.5), Inches(5.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RgbColor(0x60, 0x60, 0x60)
tf = shape.text_frame
tf.paragraphs[0].text = "[Professional working on laptop]"
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 4: The Problem
slide4 = prs.slides.add_slide(blank_layout)
set_background(slide4, MEDD_DARK)

add_text_box(slide4, 0.5, 0.5, 5, 1, "The problem we face", font_size=36, bold=True, color=WHITE)

# Problem cards
problems = [
    ("Problem 1", "Nobody enjoys role-plays, so they tend to be a filler session at sales meetings or a dreaded component of SFE.", WHITE, BLACK),
    ("Problem 2", "When managers are stretched thin, coaching is usually the first thing that is dropped.", MEDD_LIME, BLACK),
    ("Problem 3", "Most professionals love the results of coaching. They just don't love the feeling of being exposed.", MEDD_DARK, WHITE),
]

for i, (title, desc, bg, fg) in enumerate(problems):
    left = 0.5 + i * 4.2
    shape = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(4.5), Inches(4), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = fg
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = fg

# Slide 5: How We Fix It
slide5 = prs.slides.add_slide(blank_layout)
set_background(slide5, MEDD_DARK)

add_text_box(slide5, 0.5, 0.3, 6, 0.8, "How we fix it", font_size=36, bold=True, color=WHITE)
add_text_box(slide5, 0.5, 1.1, 6, 1.2, 
    "A behaviour rehearsal engine:\nbuild scenarios, practice in private,\nand make it happen in the real world.",
    font_size=18, color=WHITE)

# Solution cards
solutions = [
    ("1. Total Control", "User-generated agents: build from templates or scratch for sales, business coaching, patient simulations, OSCE-style examiners", RgbColor(0xE0, 0xE0, 0xE0)),
    ("2. Content Studio", "Custom Notebook LM style tools to generate all types of resources with a click of a button.", RgbColor(0xF0, 0xF0, 0xF0)),
    ("3. BYO Rubric", "Standardise coaching and feedback style based on the company's subscribed format.", MEDD_LIME),
    ("4. Learning Pathways", "Create structured learning pathways for product launches, clinical training, business acumen, and performance evaluation.", RgbColor(0xF5, 0xE6, 0x63)),
    ("5. Security", "• MEDD servers\n• Isolated instance\n• No LLM training\n• RAG (retrieval augmentation)", RgbColor(0x5B, 0xC0, 0xDE)),
]

for i, (title, desc, bg) in enumerate(solutions):
    left = 0.3 + i * 2.5
    shape = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(4.3), Inches(2.4), Inches(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = BLACK
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(9)
    p2.font.color.rgb = BLACK

# Slide 6: Build Any Experience
slide6 = prs.slides.add_slide(blank_layout)
set_background(slide6, MEDD_DARK)

add_text_box(slide6, 0.5, 0.3, 12, 0.8, "Build any coaching or roleplay experience", font_size=32, bold=True, color=WHITE)

experiences = [
    ("Roleplay any customer", "Create agents with specific personas, objections, and buying motives. Then set the temperature for real-world conversations."),
    ("Coach any situation", "Develop a coach for any situation using best practice templates for sales, business, HR and much more..."),
    ("Workforce Readiness", "Create learning modules and simulations for orientation, performance reviews, compliancy training."),
    ("Patient case study", "Take therapy development to the next level by adding AI patient case-study simulations to CPD presentations."),
    ("Examiner Mode", "Create an OSCE agent to assess any clinical, technical, or knowledge-based scenario with a custom rubric."),
    ("User-created Sims", "Let the team create personalised coaching, role-play, or assessment simulations tailored to their goals."),
]

for i, (title, desc) in enumerate(experiences):
    row = i // 3
    col = i % 3
    left = 0.5 + col * 4.2
    top = 1.5 + row * 2.5
    bg = WHITE if (i % 2 == 0) else MEDD_LIME
    
    shape = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(4), Inches(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = MEDD_GREEN if bg == WHITE else BLACK
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = BLACK

add_text_box(slide6, 1, 6.7, 11, 0.5, 
    "Create from templates, build from scratch, or go full service",
    font_size=18, color=WHITE, align=PP_ALIGN.CENTER)

# Slide 7: How It Works
slide7 = prs.slides.add_slide(blank_layout)
set_background(slide7, MEDD_DARK)

add_text_box(slide7, 0.5, 0.3, 6, 0.8, "How it Works", font_size=36, bold=True, color=WHITE)

steps = [
    ("1", "Register to medd to access the sim dashboard and set up your account. Select a subscriber package."),
    ("2", "Onboard the team using the onboarding tools and assign user types."),
    ("3", "Create your first agent, learning pathway, roleplay, gamified campaign or LearnDash course."),
    ("4", "Build patient simulation case studies and CPD education for therapy development activities."),
]

for i, (num, desc) in enumerate(steps):
    left = 0.5 + i * 3.2
    
    # Card
    shape = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.5), Inches(3), Inches(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MEDD_LIME
    shape.line.fill.background()
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = BLACK
    
    # Number circle
    circle = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 1.1), Inches(5.8), Inches(0.8), Inches(0.8))
    circle.fill.background()
    circle.line.color.rgb = MEDD_LIME
    circle.line.width = Pt(2)
    
    tf2 = circle.text_frame
    tf2.paragraphs[0].text = num
    tf2.paragraphs[0].font.size = Pt(24)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = MEDD_LIME
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 8: Security
slide8 = prs.slides.add_slide(blank_layout)
set_background(slide8, MEDD_DARK)

add_text_box(slide8, 0.5, 0.3, 8, 0.8, "Security and governance", font_size=36, bold=True, color=WHITE)
add_text_box(slide8, 0.5, 1, 8, 0.5, "Private by design - secure by default", font_size=20, color=MEDD_LIME)

security_sections = [
    ("Data boundary", "• No model training\n• Customer isolation (separate instance - no sharing)\n• Configurable retention"),
    ("Access & oversight", "• Role-based access (learner/assessor/admin)\n• Admin governance (scenario approval, assessor assignment)\n• Optional transcript controls (on/off, redaction)"),
    ("Audit & compliance", "• Audit logs (who did what, when)\n• Rubric versioning (which rules + content were used)\n• Exportable evidence accessible by admin/end-user"),
]

for i, (title, items) in enumerate(security_sections):
    top = 1.8 + i * 1.7
    add_text_box(slide8, 0.5, top, 3, 0.4, title, font_size=14, bold=True, color=WHITE)
    add_text_box(slide8, 0.5, top + 0.4, 7, 1.2, items, font_size=12, color=WHITE)

# Security icon placeholder
shape = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9), Inches(2), Inches(3.5), Inches(3))
shape.fill.solid()
shape.fill.fore_color.rgb = MEDD_LIME
shape.line.fill.background()
tf = shape.text_frame
tf.paragraphs[0].text = "🔒"
tf.paragraphs[0].font.size = Pt(72)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 9: Demo
slide9 = prs.slides.add_slide(blank_layout)
set_background(slide9, MEDD_DARK)

add_text_box(slide9, 2, 2.5, 9, 2, "Demo", font_size=96, bold=True, color=MEDD_LIME, align=PP_ALIGN.CENTER)

# Try to add screenshot
screenshot_path = os.path.join(os.path.dirname(__file__), 'screenshots', 'content-studio.png')
if os.path.exists(screenshot_path):
    slide9.shapes.add_picture(screenshot_path, Inches(3), Inches(4.5), width=Inches(7))

# Slide 10: Pricing
slide10 = prs.slides.add_slide(blank_layout)
set_background(slide10, MEDD_DARK)

add_text_box(slide10, 0.5, 0.3, 3, 0.8, "Pricing", font_size=36, bold=True, color=WHITE)

# Affordable badge
badge = slide10.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5), Inches(0.3), Inches(2), Inches(2))
badge.fill.solid()
badge.fill.fore_color.rgb = MEDD_LIME
badge.line.fill.background()
tf = badge.text_frame
tf.paragraphs[0].text = "AFFORDABLE"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = MEDD_GREEN
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Pricing columns
pricing = [
    ("Essential — $39/user/month", 
     "Best for individual users and small teams\n\n"
     "• 10 mins video / 60 mins audio per user/month\n"
     "• Unlimited agent types\n"
     "• Learning pathways\n"
     "• Templates + custom rubrics\n"
     "• Content Studio: 2 sessions/month\n"
     "• Email support\n\n"
     "Overage: $0.59(V)-$0.16(A)/min"),
    ("Professional — $79/user/month",
     "Best for coordinated training programs\n\n"
     "• Pooled: 25 mins video × seats/month\n"
     "• 90 mins audio × seats/month\n"
     "• Everything in Essential\n"
     "• Team analytics + gamification\n"
     "• Pooled Content Studio\n"
     "• Email + message support\n\n"
     "Overage: $0.37(V)-$0.11(A)/min"),
    ("Enterprise — Custom",
     "Best for organisations\n\n"
     "• Custom instance\n"
     "• High pooled allowances\n"
     "• SSO & advanced governance\n"
     "• Volume rates\n"
     "• Dedicated success manager\n"
     "• Phone support"),
]

for i, (title, details) in enumerate(pricing):
    left = 0.5 + i * 4.3
    add_text_box(slide10, left, 2.5, 4, 0.5, title, font_size=14, bold=True, color=WHITE)
    add_text_box(slide10, left, 3.1, 4, 4, details, font_size=10, color=WHITE)

# Slide 11: Content Studio Screenshot
slide11 = prs.slides.add_slide(blank_layout)
set_background(slide11, MEDD_DARK)

add_text_box(slide11, 0.5, 0.3, 3, 0.4, "Content Studio", font_size=16, color=MEDD_LIME)
add_text_box(slide11, 0.5, 0.7, 5, 1, "Generate learning materials in seconds", font_size=28, bold=True, color=WHITE)
add_text_box(slide11, 0.5, 2, 5, 2.5,
    "Create podcasts, flashcards, quizzes, mind maps, summaries, FAQs, briefings, and more — all from your source content.\n\n"
    "• Healthcare, Sales, Coaching content types\n"
    "• Brief, Standard, or Detailed outputs\n"
    "• Custom instructions for regeneration",
    font_size=14, color=WHITE)

# Add screenshot
screenshot_path = os.path.join(os.path.dirname(__file__), 'screenshots', 'content-studio.png')
if os.path.exists(screenshot_path):
    slide11.shapes.add_picture(screenshot_path, Inches(6.5), Inches(0.5), width=Inches(6.3))

# Slide 12: Agent Builder Screenshot
slide12 = prs.slides.add_slide(blank_layout)
set_background(slide12, MEDD_DARK)

add_text_box(slide12, 0.5, 0.3, 3, 0.4, "Agent Builder", font_size=16, color=MEDD_LIME)
add_text_box(slide12, 0.5, 0.7, 5, 1, "Create custom AI agents with full control", font_size=28, bold=True, color=WHITE)
add_text_box(slide12, 0.5, 2, 5, 2.5,
    "Define preparation context, simulation objectives, evaluation rubrics, and coaching styles — all in one intuitive interface.\n\n"
    "• Custom evaluation criteria with weighted scoring\n"
    "• 6 coaching styles: Challenger to Facilitative\n"
    "• Configurable difficulty and duration",
    font_size=14, color=WHITE)

# Add screenshot
screenshot_path = os.path.join(os.path.dirname(__file__), 'screenshots', 'agent-create.png')
if os.path.exists(screenshot_path):
    slide12.shapes.add_picture(screenshot_path, Inches(6.5), Inches(0.5), height=Inches(6.5))

# Slide 13: CTA
slide13 = prs.slides.add_slide(blank_layout)
set_background(slide13, MEDD_DARK)

add_text_box(slide13, 2, 1.5, 9, 1.5, "medd sim", font_size=48, bold=True, color=MEDD_GREEN, align=PP_ALIGN.CENTER)
add_text_box(slide13, 1.5, 3.5, 10, 1, 
    "Ready to transform how your team practices?",
    font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide13, 1.5, 4.8, 10, 0.5,
    "Turn high-stakes conversations into rehearsed performances.",
    font_size=18, color=RgbColor(0x9C, 0xA3, 0xAF), align=PP_ALIGN.CENTER)

# CTA button
button = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(5.5), Inches(3.3), Inches(0.8))
button.fill.solid()
button.fill.fore_color.rgb = MEDD_LIME
button.line.fill.background()
tf = button.text_frame
tf.paragraphs[0].text = "Start Free Trial"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = MEDD_GREEN
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

add_text_box(slide13, 2, 6.5, 9, 0.5,
    "📧 hello@medd.com.au  |  🌐 sim.medd.com.au",
    font_size=14, color=RgbColor(0x9C, 0xA3, 0xAF), align=PP_ALIGN.CENTER)

# Save
output_path = os.path.join(os.path.dirname(__file__), 'MEDD-SIM-Pitch-Deck.pptx')
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
